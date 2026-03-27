"""
LEVERAGE v3 — Python Sidecar
FastAPI service on port 5001.

Endpoints:
  POST /api/commodities          — yfinance multi-ticker fetch
  POST /api/fred                 — FRED series fetch
  POST /api/eia                  — EIA energy series fetch
  POST /api/extract/article      — Trafilatura full-text extraction
  POST /api/extract/rss          — feedparser RSS feed parse
  POST /generate/pptx            — python-pptx steerco deck (Phase 1 stub)
  POST /generate/docx            — python-docx ODD memo (Phase 1 stub)
  POST /generate/xlsx            — openpyxl Excel model (Phase 1 stub)
  GET  /health                   — health check
"""

from __future__ import annotations

import json
import logging
import os
import time
from datetime import datetime, timedelta
from typing import Any

from dotenv import load_dotenv
load_dotenv()

import httpx
import pandas as pd
import yfinance as yf
import feedparser
import trafilatura
from fastapi import FastAPI, HTTPException, Response
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger("sidecar")

app = FastAPI(title="LEVERAGE v3 Python Sidecar", version="3.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5000", "http://localhost:3000", "http://127.0.0.1:5000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
FRED_API_KEY = os.getenv("FRED_API_KEY", "")
EIA_API_KEY = os.getenv("EIA_API_KEY", "")

# All 20 commodity tickers with metadata
COMMODITY_TICKERS: dict[str, dict[str, str]] = {
    # Metals
    "GC=F":  {"name": "Gold Futures (COMEX)",          "unit": "$/oz",       "tag": "commodity_metal"},
    "HG=F":  {"name": "Copper Futures (COMEX)",         "unit": "$/lb",       "tag": "commodity_metal"},
    "ALI=F": {"name": "Aluminum Futures (COMEX)",       "unit": "$/lb",       "tag": "commodity_metal"},
    "SI=F":  {"name": "Silver Futures (COMEX)",         "unit": "$/oz",       "tag": "commodity_metal"},
    "NI=F":  {"name": "Nickel Futures",                 "unit": "$/lb",       "tag": "commodity_metal"},
    "PA=F":  {"name": "Palladium Futures",              "unit": "$/oz",       "tag": "commodity_metal"},
    "PL=F":  {"name": "Platinum Futures",               "unit": "$/oz",       "tag": "commodity_metal"},
    # Energy
    "CL=F":  {"name": "WTI Crude Oil Futures",          "unit": "$/barrel",   "tag": "commodity_energy"},
    "BZ=F":  {"name": "Brent Crude Oil Futures",        "unit": "$/barrel",   "tag": "commodity_energy"},
    "NG=F":  {"name": "Natural Gas Futures (Henry Hub)","unit": "$/MMBtu",    "tag": "commodity_energy"},
    # Ag
    "ZW=F":  {"name": "Wheat Futures (CBOT)",           "unit": "cents/bu",   "tag": "commodity_ag"},
    "ZC=F":  {"name": "Corn Futures (CBOT)",            "unit": "cents/bu",   "tag": "commodity_ag"},
    "ZS=F":  {"name": "Soybean Futures (CBOT)",         "unit": "cents/bu",   "tag": "commodity_ag"},
    "CT=F":  {"name": "Cotton Futures (ICE)",           "unit": "cents/lb",   "tag": "commodity_ag"},
    "KC=F":  {"name": "Coffee Futures (ICE)",           "unit": "cents/lb",   "tag": "commodity_ag"},
    "SB=F":  {"name": "Sugar #11 Futures (ICE)",        "unit": "cents/lb",   "tag": "commodity_ag"},
    # Lumber / Steel
    "LBS=F": {"name": "Random Length Lumber Futures",   "unit": "$/1000 board ft", "tag": "commodity_ag"},
    "HR=F":  {"name": "Hot-Rolled Coil Steel Futures",  "unit": "$/short ton","tag": "commodity_metal"},
}

# ---------------------------------------------------------------------------
# Models
# ---------------------------------------------------------------------------
class CommodityRequest(BaseModel):
    tickers: list[str] | None = None   # None = fetch all 20

class FredRequest(BaseModel):
    series_ids: list[str]              # e.g. ['CPIAUCSL', 'WPU10', 'FEDFUNDS']

class EiaRequest(BaseModel):
    series_ids: list[str]              # e.g. ['PET.RWTC.W', 'NG.RNGWHHD.W']

class ArticleExtractRequest(BaseModel):
    url: str

class RssRequest(BaseModel):
    feed_urls: list[str]
    max_items_per_feed: int = 20


# ---------------------------------------------------------------------------
# /health
# ---------------------------------------------------------------------------
@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok", "service": "leverage-v3-sidecar", "time": datetime.utcnow().isoformat()}


# ---------------------------------------------------------------------------
# POST /api/commodities  — yfinance multi-ticker
# ---------------------------------------------------------------------------
@app.post("/api/commodities")
def fetch_commodities(req: CommodityRequest) -> dict[str, Any]:
    tickers_to_fetch = req.tickers if req.tickers else list(COMMODITY_TICKERS.keys())
    log.info(f"Fetching {len(tickers_to_fetch)} commodity tickers via yfinance")

    results: list[dict[str, Any]] = []
    errors: list[str] = []

    # Fetch in one batch call — yfinance supports space-separated tickers
    ticker_str = " ".join(tickers_to_fetch)
    try:
        # Download last 2 days to compute MoM proxy (last close vs prev close)
        # For YoY we fetch 400 days
        df_recent = yf.download(
            ticker_str,
            period="2d",
            interval="1d",
            auto_adjust=True,
            progress=False,
            threads=True,
        )
        df_yoy = yf.download(
            ticker_str,
            period="400d",
            interval="1mo",
            auto_adjust=True,
            progress=False,
            threads=True,
        )
    except Exception as e:
        log.error(f"yfinance batch download failed: {e}")
        raise HTTPException(status_code=502, detail=f"yfinance error: {e}")

    close_recent = df_recent["Close"] if "Close" in df_recent.columns else df_recent
    close_yoy = df_yoy["Close"] if "Close" in df_yoy.columns else df_yoy

    for ticker in tickers_to_fetch:
        meta = COMMODITY_TICKERS.get(ticker, {"name": ticker, "unit": "", "tag": "commodity_metal"})
        try:
            # Handle single vs multi ticker DataFrame shapes
            if len(tickers_to_fetch) == 1:
                prices_recent = close_recent.dropna()
                prices_yoy = close_yoy.dropna()
            else:
                prices_recent = close_recent[ticker].dropna() if ticker in close_recent.columns else pd.Series(dtype=float)
                prices_yoy = close_yoy[ticker].dropna() if ticker in close_yoy.columns else pd.Series(dtype=float)

            if prices_recent.empty:
                errors.append(f"{ticker}: no recent data")
                continue

            latest_value = float(prices_recent.iloc[-1])
            latest_date = str(prices_recent.index[-1].date())

            # MoM: compare last two available daily closes
            mom_pct: float | None = None
            if len(prices_recent) >= 2:
                prev = float(prices_recent.iloc[-2])
                mom_pct = round((latest_value - prev) / prev * 100, 2) if prev != 0 else None

            # YoY: compare latest monthly close to ~12 months ago
            yoy_pct: float | None = None
            if len(prices_yoy) >= 13:
                val_12mo_ago = float(prices_yoy.iloc[-13])
                yoy_pct = round((latest_value - val_12mo_ago) / val_12mo_ago * 100, 2) if val_12mo_ago != 0 else None

            results.append({
                "series_id": ticker,
                "series_name": meta["name"],
                "category_tag": meta["tag"],
                "value": round(latest_value, 4),
                "unit": meta["unit"],
                "period": latest_date,
                "yoy_change_pct": yoy_pct,
                "mom_change_pct": mom_pct,
                "data_source": "yfinance",
                "ttl_hours": 1,
                "fetched_at": datetime.utcnow().isoformat(),
            })
        except Exception as e:
            errors.append(f"{ticker}: {e}")
            log.warning(f"Ticker {ticker} error: {e}")

    log.info(f"Commodity fetch complete: {len(results)} ok, {len(errors)} errors")
    return {"results": results, "errors": errors, "fetched_at": datetime.utcnow().isoformat()}


# ---------------------------------------------------------------------------
# POST /api/fred  — FRED series fetch
# ---------------------------------------------------------------------------
@app.post("/api/fred")
def fetch_fred(req: FredRequest) -> dict[str, Any]:
    if not FRED_API_KEY:
        raise HTTPException(status_code=503, detail="FRED_API_KEY not set. Add to environment.")

    FRED_META: dict[str, dict[str, str]] = {
        "CPIAUCSL": {"name": "CPI All Urban Consumers", "unit": "index", "tag": "macro"},
        "WPU10":    {"name": "PPI Metals & Metal Products", "unit": "index", "tag": "ppi"},
        "WPU06":    {"name": "PPI Chemicals & Allied Products", "unit": "index", "tag": "ppi"},
        "WPU05":    {"name": "PPI Lumber & Wood Products", "unit": "index", "tag": "ppi"},
        "WPU091":   {"name": "PPI Energy", "unit": "index", "tag": "ppi"},
        "FEDFUNDS": {"name": "Federal Funds Effective Rate", "unit": "%", "tag": "macro"},
        "GS10":     {"name": "10-Year Treasury Constant Maturity Rate", "unit": "%", "tag": "macro"},
        "ECI":      {"name": "Employment Cost Index", "unit": "index", "tag": "labor"},
    }

    results: list[dict[str, Any]] = []
    errors: list[str] = []

    async_client = httpx.Client(timeout=15)

    for series_id in req.series_ids:
        try:
            url = "https://api.stlouisfed.org/fred/series/observations"
            params = {
                "series_id": series_id,
                "api_key": FRED_API_KEY,
                "file_type": "json",
                "sort_order": "desc",
                "limit": 14,          # last 14 observations for YoY/MoM calc
            }
            resp = async_client.get(url, params=params)
            resp.raise_for_status()
            data = resp.json()
            obs = [o for o in data.get("observations", []) if o["value"] != "."]

            if not obs:
                errors.append(f"{series_id}: no observations")
                continue

            latest_val = float(obs[0]["value"])
            latest_period = obs[0]["date"][:7]  # YYYY-MM

            mom_pct: float | None = None
            if len(obs) >= 2:
                prev = float(obs[1]["value"])
                mom_pct = round((latest_val - prev) / prev * 100, 2) if prev != 0 else None

            yoy_pct: float | None = None
            if len(obs) >= 13:
                prev_yr = float(obs[12]["value"])
                yoy_pct = round((latest_val - prev_yr) / prev_yr * 100, 2) if prev_yr != 0 else None

            meta = FRED_META.get(series_id, {"name": series_id, "unit": "index", "tag": "macro"})
            results.append({
                "series_id": series_id,
                "series_name": meta["name"],
                "category_tag": meta["tag"],
                "value": latest_val,
                "unit": meta["unit"],
                "period": latest_period,
                "yoy_change_pct": yoy_pct,
                "mom_change_pct": mom_pct,
                "data_source": "fred",
                "ttl_hours": 24,
                "fetched_at": datetime.utcnow().isoformat(),
                "raw_json": json.dumps(obs[:14]),
            })
        except Exception as e:
            errors.append(f"{series_id}: {e}")
            log.warning(f"FRED {series_id} error: {e}")

    async_client.close()
    return {"results": results, "errors": errors, "fetched_at": datetime.utcnow().isoformat()}


# ---------------------------------------------------------------------------
# POST /api/eia  — EIA energy series
# ---------------------------------------------------------------------------
@app.post("/api/eia")
def fetch_eia(req: EiaRequest) -> dict[str, Any]:
    if not EIA_API_KEY:
        raise HTTPException(status_code=503, detail="EIA_API_KEY not set. Add to environment.")

    EIA_META: dict[str, dict[str, str]] = {
        "PET.RWTC.W":    {"name": "WTI Crude Oil Weekly Spot Price", "unit": "$/barrel", "tag": "commodity_energy"},
        "NG.RNGWHHD.W":  {"name": "Henry Hub Natural Gas Weekly Spot", "unit": "$/MMBtu", "tag": "commodity_energy"},
    }

    results: list[dict[str, Any]] = []
    errors: list[str] = []
    client = httpx.Client(timeout=15)

    for series_id in req.series_ids:
        try:
            # EIA v2 API
            url = f"https://api.eia.gov/v2/seriesid/{series_id}"
            params = {"api_key": EIA_API_KEY, "data[0]": "value", "sort[0][column]": "period", "sort[0][direction]": "desc", "length": 56}
            resp = client.get(url, params=params)
            resp.raise_for_status()
            data = resp.json()
            obs = data.get("response", {}).get("data", [])

            if not obs:
                errors.append(f"{series_id}: no data")
                continue

            latest_val = float(obs[0]["value"])
            latest_period = obs[0]["period"]

            # WoW (week over week)
            mom_pct: float | None = None
            if len(obs) >= 2:
                prev = float(obs[1]["value"])
                mom_pct = round((latest_val - prev) / prev * 100, 2) if prev != 0 else None

            # YoY ~52 weeks
            yoy_pct: float | None = None
            if len(obs) >= 53:
                prev_yr = float(obs[52]["value"])
                yoy_pct = round((latest_val - prev_yr) / prev_yr * 100, 2) if prev_yr != 0 else None

            meta = EIA_META.get(series_id, {"name": series_id, "unit": "value", "tag": "commodity_energy"})
            results.append({
                "series_id": series_id,
                "series_name": meta["name"],
                "category_tag": meta["tag"],
                "value": latest_val,
                "unit": meta["unit"],
                "period": latest_period,
                "yoy_change_pct": yoy_pct,
                "mom_change_pct": mom_pct,
                "data_source": "eia",
                "ttl_hours": 12,
                "fetched_at": datetime.utcnow().isoformat(),
            })
        except Exception as e:
            errors.append(f"{series_id}: {e}")
            log.warning(f"EIA {series_id} error: {e}")

    client.close()
    return {"results": results, "errors": errors, "fetched_at": datetime.utcnow().isoformat()}


# ---------------------------------------------------------------------------
# POST /api/extract/article  — Trafilatura full-text extraction
# ---------------------------------------------------------------------------
@app.post("/api/extract/article")
def extract_article(req: ArticleExtractRequest) -> dict[str, Any]:
    url = req.url
    log.info(f"Extracting article: {url}")

    # Attempt 1: Direct fetch
    downloaded = trafilatura.fetch_url(url)
    text = trafilatura.extract(downloaded) if downloaded else None
    confidence = "high"

    # Attempt 2: archive.ph fallback for paywalled content
    if not text or len(text) < 200:
        archive_url = f"https://archive.ph/{url}"
        try:
            downloaded_archive = trafilatura.fetch_url(archive_url)
            text_archive = trafilatura.extract(downloaded_archive) if downloaded_archive else None
            if text_archive and len(text_archive) >= 200:
                text = text_archive
                confidence = "medium"
            else:
                confidence = "low"
                text = text or ""
        except Exception:
            confidence = "low"
            text = text or ""

    return {
        "url": url,
        "text": text or "",
        "confidence": confidence,
        "char_count": len(text or ""),
        "extracted_at": datetime.utcnow().isoformat(),
    }


# ---------------------------------------------------------------------------
# POST /api/extract/rss  — feedparser RSS feed parse
# ---------------------------------------------------------------------------
TRADE_FEEDS = [
    "https://www.supplychaindive.com/feeds/news/",
    "https://spendmatters.com/feed/",
    "https://www.freightwaves.com/news/feed",
    "https://feeds.reuters.com/reuters/businessNews",
]

@app.post("/api/extract/rss")
def extract_rss(req: RssRequest) -> dict[str, Any]:
    feed_urls = req.feed_urls if req.feed_urls else TRADE_FEEDS
    all_items: list[dict[str, Any]] = []

    for feed_url in feed_urls:
        try:
            feed = feedparser.parse(feed_url)
            for entry in feed.entries[:req.max_items_per_feed]:
                all_items.append({
                    "feed_url": feed_url,
                    "title": entry.get("title", ""),
                    "summary": entry.get("summary", ""),
                    "link": entry.get("link", ""),
                    "published": entry.get("published", ""),
                    "tags": [t.get("term", "") for t in entry.get("tags", [])],
                })
        except Exception as e:
            log.warning(f"RSS feed {feed_url} error: {e}")

    return {
        "items": all_items,
        "item_count": len(all_items),
        "fetched_at": datetime.utcnow().isoformat(),
    }


# ---------------------------------------------------------------------------
# POST /generate/pptx  — Steerco deck (P1-09)
# ---------------------------------------------------------------------------
class PptxRequest(BaseModel):
    engagement: dict[str, Any]
    slides: dict[str, Any]
    branding: dict[str, Any] = {}


@app.post("/generate/pptx")
def generate_pptx(req: PptxRequest) -> Response:
    """
    Generate a branded steerco PPTX using python-pptx.
    Returns the file as application/octet-stream.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    import io

    def hex_to_rgb(hex_color: str) -> RGBColor:
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    primary = hex_to_rgb(req.branding.get("primary_color", "#003366"))
    secondary = hex_to_rgb(req.branding.get("secondary_color", "#0066CC"))
    white = RGBColor(0xFF, 0xFF, 0xFF)
    light_gray = RGBColor(0xF5, 0xF5, 0xF5)
    dark_gray = RGBColor(0x33, 0x33, 0x33)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    def add_slide() -> Any:
        return prs.slides.add_slide(blank_layout)

    def add_rect(slide: Any, left: float, top: float, width: float, height: float, fill_color: RGBColor | None = None) -> Any:
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        shape.line.fill.background()
        return shape

    def add_text_box(slide: Any, text: str, left: float, top: float, width: float, height: float,
                     font_size: int = 14, bold: bool = False, color: RGBColor | None = None,
                     align: Any = PP_ALIGN.LEFT, wrap: bool = True) -> Any:
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def slide_header(slide: Any, title: str, subtitle: str = "") -> None:
        add_rect(slide, 0, 0, 13.33, 1.2, primary)
        add_text_box(slide, title, 0.3, 0.15, 10, 0.6, font_size=24, bold=True, color=white)
        if subtitle:
            add_text_box(slide, subtitle, 0.3, 0.75, 10, 0.35, font_size=12, color=white)
        # Footer bar
        add_rect(slide, 0, 7.1, 13.33, 0.4, primary)
        footer_text = req.branding.get("header_text", "CONFIDENTIAL — A&M PEPI")
        add_text_box(slide, f"{footer_text}  |  {req.engagement.get('portfolio_company', '')}",
                     0.3, 7.12, 10, 0.3, font_size=8, color=white)

    company = req.engagement.get("portfolio_company", "Portfolio Company")
    sponsor = req.engagement.get("pe_sponsor", "PE Sponsor")

    # -------------------------------------------------------------------
    # Slide 1: Title
    # -------------------------------------------------------------------
    slide = add_slide()
    add_rect(slide, 0, 0, 13.33, 7.5, primary)
    add_rect(slide, 0.3, 2.8, 12.73, 0.05, secondary)
    add_text_box(slide, company, 0.5, 1.5, 12, 1.0, font_size=36, bold=True, color=white, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Procurement Opportunity Assessment", 0.5, 2.9, 12, 0.6, font_size=20, color=white, align=PP_ALIGN.CENTER)
    add_text_box(slide, f"Prepared for {sponsor}  |  {req.engagement.get('name', '')}",
                 0.5, 3.6, 12, 0.4, font_size=14, color=white, align=PP_ALIGN.CENTER)
    add_text_box(slide, f"A&M PEPI  |  {req.engagement.get('industry', '')}  |  Confidential",
                 0.5, 6.8, 12, 0.4, font_size=10, color=white, align=PP_ALIGN.CENTER)

    # -------------------------------------------------------------------
    # Slide 2: Situation
    # -------------------------------------------------------------------
    situation = req.slides.get("situation", {})
    slide = add_slide()
    slide_header(slide, "Situation", f"Total Addressable Spend: {situation.get('total_spend', '—')}")
    add_text_box(slide, situation.get("narrative", ""), 0.4, 1.35, 8.5, 2.0, font_size=13, color=dark_gray, wrap=True)
    # Spend breakdown mini-table
    breakdown = situation.get("spend_breakdown", [])
    if breakdown:
        add_text_box(slide, "Spend by Category", 9.3, 1.35, 3.6, 0.4, font_size=11, bold=True, color=primary)
        for idx, row in enumerate(breakdown[:6]):
            y = 1.8 + idx * 0.37
            bg = light_gray if idx % 2 == 0 else white
            add_rect(slide, 9.3, y, 3.6, 0.35, bg)
            add_text_box(slide, str(row.get("category", "")), 9.4, y + 0.04, 2.2, 0.28, font_size=10, color=dark_gray)
            add_text_box(slide, str(row.get("spend", "")), 11.4, y + 0.04, 1.4, 0.28, font_size=10, bold=True,
                         color=primary, align=PP_ALIGN.RIGHT)

    # -------------------------------------------------------------------
    # Slide 3: Savings Initiative Pipeline
    # -------------------------------------------------------------------
    inits = req.slides.get("initiatives", {})
    slide = add_slide()
    slide_header(slide, "Savings Initiative Pipeline",
                 f"Total: {inits.get('pipeline_total', '—')}  |  Risk-Adjusted: {inits.get('risk_adjusted', '—')}")

    # KPI cards
    kpis = [
        ("Pipeline", inits.get("pipeline_total", "—"), secondary),
        ("Risk-Adjusted", inits.get("risk_adjusted", "—"), primary),
        ("Quick Wins", inits.get("quick_win", "—"), RGBColor(0x00, 0x88, 0x44)),
    ]
    for i, (label, value, color) in enumerate(kpis):
        x = 0.4 + i * 4.0
        add_rect(slide, x, 1.3, 3.5, 1.1, color)
        add_text_box(slide, label, x + 0.1, 1.35, 3.3, 0.4, font_size=11, color=white)
        add_text_box(slide, str(value), x + 0.1, 1.65, 3.3, 0.6, font_size=22, bold=True, color=white)

    # Initiatives table
    items = inits.get("items", [])
    headers = ["Initiative", "Lever", "Phase", "Target", "Risk-Adj", "Prob"]
    col_widths = [3.8, 2.0, 1.6, 1.3, 1.3, 0.9]
    col_starts = [0.4]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    # Header row
    add_rect(slide, 0.4, 2.55, sum(col_widths), 0.35, primary)
    for j, (h, x) in enumerate(zip(headers, col_starts)):
        add_text_box(slide, h, x + 0.05, 2.57, col_widths[j] - 0.1, 0.3, font_size=9, bold=True, color=white)

    for row_i, item in enumerate(items[:7]):
        y = 2.95 + row_i * 0.35
        bg = light_gray if row_i % 2 == 0 else white
        add_rect(slide, 0.4, y, sum(col_widths), 0.33, bg)
        row_vals = [
            str(item.get("name", ""))[:42],
            str(item.get("lever", "") or "")[:20],
            str(item.get("phase", "") or "").replace("_", " "),
            str(item.get("target", "—")),
            str(item.get("risk_adjusted", "—")),
            str(item.get("probability", "—")),
        ]
        for j, (val, x) in enumerate(zip(row_vals, col_starts)):
            add_text_box(slide, val, x + 0.05, y + 0.04, col_widths[j] - 0.1, 0.28,
                         font_size=9, color=dark_gray)

    # -------------------------------------------------------------------
    # Slide 4: Kraljic Matrix (text-based quadrant summary)
    # -------------------------------------------------------------------
    kraljic = req.slides.get("kraljic", {})
    slide = add_slide()
    slide_header(slide, "Category Strategy — Kraljic Matrix", "Procurement risk vs. profit impact positioning")

    quadrants = {"Strategic": [], "Leverage": [], "Bottleneck": [], "Non-critical": []}
    for cat in kraljic.get("categories", []):
        q = str(cat.get("quadrant", "Non-critical"))
        if q in quadrants:
            quadrants[q].append(str(cat.get("name", "")))

    quad_config = [
        ("Strategic", 0.4, 1.3, secondary, "High risk, high impact\nPartner deeply, multi-year contracts"),
        ("Leverage", 6.9, 1.3, RGBColor(0x00, 0x88, 0x44), "Low risk, high impact\nCompetitive bidding, volume consolidation"),
        ("Bottleneck", 0.4, 4.2, RGBColor(0xCC, 0x55, 0x00), "High risk, low impact\nSecure supply, dual-source"),
        ("Non-critical", 6.9, 4.2, RGBColor(0x66, 0x66, 0x66), "Low risk, low impact\nAutomate, simplify, tail spend programs"),
    ]
    for quad_name, x, y, color, desc in quad_config:
        add_rect(slide, x, y, 6.0, 2.7, color)
        add_text_box(slide, quad_name, x + 0.15, y + 0.1, 5.7, 0.45, font_size=14, bold=True, color=white)
        add_text_box(slide, desc, x + 0.15, y + 0.55, 5.7, 0.5, font_size=9, color=white, wrap=True)
        cats = quadrants.get(quad_name, [])
        cats_text = "\n".join(f"• {c}" for c in cats[:4])
        if cats:
            add_text_box(slide, cats_text, x + 0.15, y + 1.1, 5.7, 1.4, font_size=9, color=white, wrap=True)

    # Axis labels
    add_text_box(slide, "← Low Supply Risk      High Supply Risk →", 0.4, 7.0, 12.5, 0.3,
                 font_size=9, color=dark_gray, align=PP_ALIGN.CENTER)
    add_text_box(slide, "↑ High\nProfit\nImpact\n\n↓ Low", 12.95, 1.3, 0.3, 5.6, font_size=8, color=dark_gray)

    # -------------------------------------------------------------------
    # Slide 5: EBITDA Bridge
    # -------------------------------------------------------------------
    bridge = req.slides.get("ebitda_bridge", {})
    slide = add_slide()
    slide_header(slide, "EBITDA Bridge — Savings Impact", "Risk-adjusted savings contribution by phase")

    phases = bridge.get("phases", [])
    total_savings = bridge.get("savings_impact", 0)
    bar_colors = [RGBColor(0x00, 0x88, 0x44), secondary, primary]
    max_val = max((p.get("value", 0) for p in phases), default=1) or 1
    bar_area_height = 3.5
    bar_width = 2.5
    bar_start_x = 1.5

    add_text_box(slide, f"Total Risk-Adjusted Savings: ${total_savings/1e6:.1f}M" if total_savings >= 1e6
                 else f"Total Risk-Adjusted Savings: ${total_savings:,.0f}",
                 0.4, 1.3, 12, 0.5, font_size=16, bold=True, color=primary)

    for i, phase in enumerate(phases[:3]):
        val = phase.get("value", 0) or 0
        bar_h = (val / max_val) * bar_area_height if max_val else 0
        x = bar_start_x + i * 3.5
        # Bar
        add_rect(slide, x, 1.9 + (bar_area_height - bar_h), bar_width, bar_h, bar_colors[i % 3])
        # Value label
        val_str = f"${val/1e6:.1f}M" if val >= 1e6 else f"${val:,.0f}"
        add_text_box(slide, val_str, x, 1.75, bar_width, 0.35, font_size=13, bold=True,
                     color=bar_colors[i % 3], align=PP_ALIGN.CENTER)
        # Phase label
        add_text_box(slide, str(phase.get("label", "")), x - 0.1, 5.5, bar_width + 0.2, 0.5,
                     font_size=10, color=dark_gray, align=PP_ALIGN.CENTER, wrap=True)

    # -------------------------------------------------------------------
    # Slide 6: 100-Day Roadmap
    # -------------------------------------------------------------------
    roadmap = req.slides.get("hundred_day_roadmap", {})
    slide = add_slide()
    slide_header(slide, "100-Day Implementation Roadmap", "Priority initiatives and sequencing")

    phases_data = roadmap.get("phases", [])
    phase_colors = [RGBColor(0x00, 0x66, 0xCC), RGBColor(0x00, 0x88, 0x44), RGBColor(0xCC, 0x55, 0x00)]
    for i, phase in enumerate(phases_data[:3]):
        x = 0.4 + i * 4.3
        add_rect(slide, x, 1.3, 4.0, 0.55, phase_colors[i % 3])
        add_text_box(slide, str(phase.get("label", "")), x + 0.1, 1.35, 3.8, 0.45,
                     font_size=12, bold=True, color=white, wrap=True)
        activities = phase.get("activities", [])
        for j, act in enumerate(activities[:4]):
            y = 1.95 + j * 0.55
            add_rect(slide, x, y, 4.0, 0.48, light_gray if j % 2 == 0 else white)
            add_text_box(slide, f"• {act}", x + 0.1, y + 0.06, 3.8, 0.38, font_size=10,
                         color=dark_gray, wrap=True)

    # -------------------------------------------------------------------
    # Slide 7: Risks
    # -------------------------------------------------------------------
    risks = req.slides.get("risks", {})
    slide = add_slide()
    slide_header(slide, "Key Risks & Mitigations")

    narrative = risks.get("narrative", "")
    add_text_box(slide, narrative, 0.4, 1.35, 12.5, 2.0, font_size=12, color=dark_gray, wrap=True)

    risk_items = risks.get("items", [])
    severity_colors = {"Critical": RGBColor(0xCC, 0x00, 0x00), "High": RGBColor(0xCC, 0x55, 0x00),
                       "Medium": secondary, "Low": RGBColor(0x00, 0x88, 0x44)}
    for i, risk in enumerate(risk_items[:5]):
        y = 3.5 + i * 0.55
        sev = str(risk.get("severity", "Medium"))
        sev_color = severity_colors.get(sev, secondary)
        add_rect(slide, 0.4, y, 0.8, 0.45, sev_color)
        add_text_box(slide, sev, 0.42, y + 0.08, 0.76, 0.32, font_size=9, bold=True, color=white, align=PP_ALIGN.CENTER)
        add_text_box(slide, str(risk.get("risk", "")), 1.3, y + 0.06, 11.5, 0.35, font_size=11, color=dark_gray)

    # -------------------------------------------------------------------
    # Slide 8: Next Steps
    # -------------------------------------------------------------------
    next_steps = req.slides.get("next_steps", {})
    slide = add_slide()
    slide_header(slide, "Next Steps")

    items = next_steps.get("items", [])
    step_colors = [primary, secondary, RGBColor(0x00, 0x88, 0x44), RGBColor(0xCC, 0x55, 0x00), RGBColor(0x66, 0x00, 0x99)]
    for i, step in enumerate(items[:5]):
        y = 1.5 + i * 0.95
        color = step_colors[i % len(step_colors)]
        add_rect(slide, 0.4, y, 0.7, 0.7, color)
        add_text_box(slide, str(i + 1), 0.42, y + 0.1, 0.66, 0.5, font_size=20, bold=True,
                     color=white, align=PP_ALIGN.CENTER)
        add_text_box(slide, str(step), 1.2, y + 0.1, 11.7, 0.6, font_size=12, color=dark_gray, wrap=True)

    # Save to buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=steerco.pptx"},
    )


# ---------------------------------------------------------------------------
# POST /generate/docx  — ODD Memo (P1-10)
# ---------------------------------------------------------------------------
class DocxRequest(BaseModel):
    engagement: dict[str, Any]
    summary_metrics: dict[str, Any]
    sections: dict[str, Any]
    data_tables: dict[str, Any]
    branding: dict[str, Any] = {}


@app.post("/generate/docx")
def generate_docx(req: DocxRequest) -> Response:
    """
    Generate a branded ODD memo DOCX using python-docx.
    Returns the file as application/octet-stream.
    """
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor as DocxRGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import io

    def hex_to_docx_rgb(hex_color: str) -> DocxRGB:
        h = hex_color.lstrip("#")
        return DocxRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    primary_hex = req.branding.get("primary_color", "#003366").lstrip("#")
    primary_rgb = hex_to_docx_rgb(req.branding.get("primary_color", "#003366"))

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    def set_para_color(para: Any, hex_color: str) -> None:
        for run in para.runs:
            run.font.color.rgb = hex_to_docx_rgb(hex_color)

    def add_heading(text: str, level: int = 1, color_hex: str = "#003366") -> Any:
        para = doc.add_heading(text, level=level)
        set_para_color(para, color_hex)
        return para

    def add_body(text: str, indent: float = 0) -> Any:
        para = doc.add_paragraph(text)
        para.paragraph_format.space_after = Pt(6)
        if indent:
            para.paragraph_format.left_indent = Inches(indent)
        return para

    def add_table(headers: list[str], rows: list[list[str]], col_widths: list[float] | None = None) -> Any:
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = "Table Grid"
        table.alignment = WD_TABLE_ALIGNMENT.LEFT
        hdr_cells = table.rows[0].cells
        for i, h in enumerate(headers):
            hdr_cells[i].text = h
            run = hdr_cells[i].paragraphs[0].runs[0] if hdr_cells[i].paragraphs[0].runs else hdr_cells[i].paragraphs[0].add_run(h)
            run.font.bold = True
            run.font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)
            # Set cell background
            tc = hdr_cells[i]._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:fill"), primary_hex)
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:val"), "clear")
            tcPr.append(shd)
        for row_data in rows:
            row_cells = table.add_row().cells
            for i, cell_text in enumerate(row_data):
                row_cells[i].text = str(cell_text)
                row_cells[i].paragraphs[0].runs[0].font.size = Pt(9) if row_cells[i].paragraphs[0].runs else None
        return table

    company = req.engagement.get("portfolio_company", "Portfolio Company")
    sponsor = req.engagement.get("pe_sponsor", "")
    header_text = req.branding.get("header_text", "CONFIDENTIAL — A&M PEPI")

    # -------------------------------------------------------------------
    # Cover / Header
    # -------------------------------------------------------------------
    title = doc.add_heading("", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(f"Procurement Operational Due Diligence")
    run.font.size = Pt(22)
    run.font.color.rgb = primary_rgb

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = subtitle.add_run(company)
    r.font.size = Pt(16)
    r.font.bold = True

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"Prepared for: {sponsor}  |  Date: {req.engagement.get('generated_date', '')}\n")
    meta.add_run(header_text).font.size = Pt(9)

    doc.add_paragraph()

    # Summary metrics table
    m = req.summary_metrics
    add_heading("Key Metrics", level=2)
    metrics_rows = [
        ["Total Addressable Spend", str(m.get("total_spend", "—"))],
        ["Savings Pipeline (Gross)", str(m.get("pipeline_total", "—"))],
        ["Risk-Adjusted Savings", str(m.get("risk_adjusted", "—"))],
        ["Quick Wins (0–90 days)", str(m.get("quick_wins", "—"))],
        ["Number of Initiatives", str(m.get("initiative_count", "—"))],
        ["Procurement Maturity Score", f"{m.get('maturity_score', '—')} / 5.0"],
        ["Tariff Exposure", str(m.get("tariff_exposure", "—"))],
        ["Contracts Expiring <90 Days", str(m.get("contracts_expiring_90d", "—"))],
    ]
    add_table(["Metric", "Value"], metrics_rows, [3.5, 2.5])
    doc.add_paragraph()

    # -------------------------------------------------------------------
    # Sections
    # -------------------------------------------------------------------
    sections_order = [
        ("1. Executive Summary", "executive_summary"),
        ("2. Methodology", "methodology"),
        ("3. Spend Analysis Findings", "spend_findings"),
        ("4. Savings Initiative Pipeline", "initiative_pipeline"),
        ("5. Risk Assessment", "risk_matrix"),
        ("6. Implementation Roadmap", "implementation_roadmap"),
    ]

    for title_text, key in sections_order:
        add_heading(title_text, level=1)
        content = req.sections.get(key, "")
        if content:
            # Split on double newlines for paragraphs
            for para_text in content.split("\n\n"):
                para_text = para_text.strip()
                if para_text:
                    add_body(para_text)
        doc.add_paragraph()

    # -------------------------------------------------------------------
    # Data Appendix
    # -------------------------------------------------------------------
    doc.add_page_break()
    add_heading("Appendix A: Spend by Category", level=1)
    spend_rows = req.data_tables.get("spend_by_category", [])
    if spend_rows:
        add_table(
            ["Category", "Annual Spend", "Suppliers", "% of Total"],
            [[str(r.get("category", "")), str(r.get("spend", "")),
              str(r.get("suppliers", "")), str(r.get("pct_of_total", ""))]
             for r in spend_rows],
        )

    doc.add_paragraph()
    add_heading("Appendix B: Savings Initiative Pipeline", level=1)
    init_rows = req.data_tables.get("top_initiatives", [])
    if init_rows:
        add_table(
            ["Initiative", "Lever", "Phase", "Target", "Risk-Adj.", "Confidence"],
            [[str(r.get("name", "")), str(r.get("lever", "")), str(r.get("phase", "")),
              str(r.get("target", "")), str(r.get("risk_adjusted", "")), str(r.get("confidence", ""))]
             for r in init_rows],
        )

    doc.add_paragraph()
    add_heading("Appendix C: Procurement Maturity Assessment", level=1)
    mat_rows = req.data_tables.get("maturity_dimensions", [])
    if mat_rows:
        add_table(
            ["Dimension", "Score (1–5)", "Gap to Next Level"],
            [[str(r.get("dimension", "")), str(r.get("score", "")), str(r.get("gap", "") or "—")]
             for r in mat_rows],
        )

    # Save
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=ODD_memo.docx"},
    )


# ---------------------------------------------------------------------------
# POST /generate/xlsx  — Excel model (P1-11)
# ---------------------------------------------------------------------------
class XlsxRequest(BaseModel):
    engagement: dict[str, Any]
    tabs: dict[str, Any]


@app.post("/generate/xlsx")
def generate_xlsx(req: XlsxRequest) -> Response:
    """Generate a multi-tab Excel model using openpyxl."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    import io

    PRIMARY_HEX = "003366"
    SECONDARY_HEX = "0066CC"
    HEADER_FONT = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
    HEADER_FILL = PatternFill("solid", fgColor=PRIMARY_HEX)
    SUBHEADER_FILL = PatternFill("solid", fgColor=SECONDARY_HEX)
    TITLE_FONT = Font(name="Calibri", bold=True, color=PRIMARY_HEX, size=14)
    NORMAL_FONT = Font(name="Calibri", size=10)
    ALT_FILL = PatternFill("solid", fgColor="F5F8FF")
    thin = Side(style="thin", color="CCCCCC")
    BORDER = Border(left=thin, right=thin, top=thin, bottom=thin)

    def style_header(ws: Any, row: int, max_col: int) -> None:
        for col in range(1, max_col + 1):
            cell = ws.cell(row=row, column=col)
            cell.font = HEADER_FONT
            cell.fill = HEADER_FILL
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = BORDER

    def style_data_row(ws: Any, row: int, max_col: int, alt: bool = False) -> None:
        for col in range(1, max_col + 1):
            cell = ws.cell(row=row, column=col)
            cell.font = NORMAL_FONT
            if alt:
                cell.fill = ALT_FILL
            cell.border = BORDER
            cell.alignment = Alignment(wrap_text=False)

    wb = Workbook()

    # -------------------------------------------------------------------
    # Tab 1: Summary
    # -------------------------------------------------------------------
    ws = wb.active
    ws.title = "Summary"
    summary = req.tabs.get("summary", {})

    ws["A1"] = f"Procurement Model — {req.engagement.get('portfolio_company', '')}"
    ws["A1"].font = TITLE_FONT
    ws.merge_cells("A1:D1")
    ws["A2"] = f"PE Sponsor: {req.engagement.get('pe_sponsor', '')}  |  Discount Rate: {(req.engagement.get('discount_rate') or 0.10) * 100:.1f}%"
    ws["A2"].font = Font(name="Calibri", italic=True, color="666666", size=10)
    ws.merge_cells("A2:D2")

    metrics = [
        ("Total Addressable Spend", summary.get("total_spend")),
        ("Gross Pipeline", summary.get("pipeline_total")),
        ("Risk-Adjusted Pipeline", summary.get("risk_adjusted")),
        ("Quick Wins (0–90d)", summary.get("quick_wins")),
        ("Medium Term (90–180d)", summary.get("by_phase", {}).get("medium_term")),
        ("Long Term (180d+)", summary.get("by_phase", {}).get("long_term")),
        ("Initiative Count", summary.get("initiative_count")),
    ]
    ws["A4"] = "Metric"
    ws["B4"] = "Value"
    style_header(ws, 4, 2)
    for i, (label, val) in enumerate(metrics):
        r = 5 + i
        ws.cell(r, 1, label).font = NORMAL_FONT
        ws.cell(r, 2, val if val is not None else "—").font = NORMAL_FONT
        if isinstance(val, (int, float)) and val and val > 1000:
            ws.cell(r, 2).number_format = '"$"#,##0'
        style_data_row(ws, r, 2, alt=i % 2 == 0)
    ws.column_dimensions["A"].width = 30
    ws.column_dimensions["B"].width = 18

    # -------------------------------------------------------------------
    # Tab 2: Initiative Pipeline
    # -------------------------------------------------------------------
    ws2 = wb.create_sheet("Initiative Pipeline")
    headers2 = ["Initiative", "Category", "Lever Type", "Phase", "Status",
                "Target ($)", "Risk-Adjusted ($)", "Probability", "Confidence", "At Risk?"]
    for col, h in enumerate(headers2, 1):
        ws2.cell(1, col, h)
    style_header(ws2, 1, len(headers2))

    inits = req.tabs.get("initiative_pipeline", [])
    for i, row in enumerate(inits):
        r = i + 2
        vals = [
            row.get("name"), row.get("category"), row.get("lever"),
            str(row.get("phase") or "").replace("_", " "),
            row.get("status"), row.get("target"), row.get("risk_adjusted"),
            row.get("probability"), row.get("confidence"),
            "Yes" if row.get("is_at_risk") else "No",
        ]
        for col, v in enumerate(vals, 1):
            cell = ws2.cell(r, col, v)
            if col in (6, 7) and isinstance(v, (int, float)):
                cell.number_format = '"$"#,##0'
            if col == 8 and isinstance(v, float):
                cell.number_format = "0%"
        style_data_row(ws2, r, len(headers2), alt=i % 2 == 0)

    col_widths2 = [35, 18, 18, 14, 14, 14, 16, 12, 12, 10]
    for col, w in enumerate(col_widths2, 1):
        ws2.column_dimensions[get_column_letter(col)].width = w
    ws2.freeze_panes = "A2"
    ws2.auto_filter.ref = f"A1:{get_column_letter(len(headers2))}1"

    # -------------------------------------------------------------------
    # Tab 3: Spend Analysis
    # -------------------------------------------------------------------
    ws3 = wb.create_sheet("Spend Analysis")
    headers3 = ["Category", "Annual Spend ($)", "Supplier Count"]
    for col, h in enumerate(headers3, 1):
        ws3.cell(1, col, h)
    style_header(ws3, 1, len(headers3))

    spend_rows = req.tabs.get("spend_analysis", [])
    for i, row in enumerate(spend_rows):
        r = i + 2
        ws3.cell(r, 1, row.get("category") or "Uncategorized")
        spend_cell = ws3.cell(r, 2, row.get("spend"))
        spend_cell.number_format = '"$"#,##0'
        ws3.cell(r, 3, row.get("suppliers"))
        style_data_row(ws3, r, 3, alt=i % 2 == 0)
    ws3.column_dimensions["A"].width = 30
    ws3.column_dimensions["B"].width = 18
    ws3.column_dimensions["C"].width = 16
    ws3.freeze_panes = "A2"

    # -------------------------------------------------------------------
    # Tab 4: Assumptions & Market Data
    # -------------------------------------------------------------------
    ws4 = wb.create_sheet("Assumptions")
    ws4["A1"] = "Model Assumptions"
    ws4["A1"].font = TITLE_FONT
    ws4["A3"] = "Discount Rate"
    ws4["B3"] = req.tabs.get("assumptions", {}).get("discount_rate", 0.10)
    ws4["B3"].number_format = "0.0%"

    mkt = req.tabs.get("assumptions", {}).get("market_data", [])
    if mkt:
        ws4["A5"] = "Market Data (from cache)"
        ws4["A5"].font = Font(name="Calibri", bold=True, color=PRIMARY_HEX)
        mkt_headers = ["Series", "Value", "Unit", "YoY %", "Period"]
        for col, h in enumerate(mkt_headers, 1):
            ws4.cell(6, col, h)
        style_header(ws4, 6, len(mkt_headers))
        for i, m in enumerate(mkt):
            r = 7 + i
            ws4.cell(r, 1, m.get("series"))
            ws4.cell(r, 2, m.get("value"))
            ws4.cell(r, 3, m.get("unit"))
            yoy = m.get("yoy_pct")
            ws4.cell(r, 4, f"{yoy:+.1f}%" if isinstance(yoy, (int, float)) else "—")
            ws4.cell(r, 5, m.get("period"))
            style_data_row(ws4, r, 5, alt=i % 2 == 0)
        for col, w in enumerate([35, 12, 14, 10, 14], 1):
            ws4.column_dimensions[get_column_letter(col)].width = w

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=model.xlsx"},
    )


# ---------------------------------------------------------------------------
# Search models
# ---------------------------------------------------------------------------
class WebSearchRequest(BaseModel):
    query: str
    max_results: int = 10
    search_type: str = "web"   # "web" | "news"


class NewsRequest(BaseModel):
    query: str
    max_results: int = 20
    language: str = "en"


# ---------------------------------------------------------------------------
# POST /api/search/ddgs  — DuckDuckGo web search (DDGS fallback)
# ---------------------------------------------------------------------------
@app.post("/api/search/ddgs")
def search_ddgs(req: WebSearchRequest) -> dict[str, Any]:
    """
    DDGS web search — used as fallback when SearXNG is unavailable.
    Returns list of {title, url, snippet} results.
    """
    try:
        from duckduckgo_search import DDGS
    except ImportError:
        raise HTTPException(status_code=503, detail="ddgs not installed. Run: pip install ddgs")

    results: list[dict[str, Any]] = []
    try:
        with DDGS() as ddgs:
            if req.search_type == "news":
                raw = ddgs.news(req.query, max_results=req.max_results)
                for r in raw:
                    results.append({
                        "title": r.get("title", ""),
                        "url": r.get("url", ""),
                        "snippet": r.get("body", ""),
                        "published": r.get("date", ""),
                        "source": r.get("source", ""),
                    })
            else:
                raw = ddgs.text(req.query, max_results=req.max_results)
                for r in raw:
                    results.append({
                        "title": r.get("title", ""),
                        "url": r.get("href", ""),
                        "snippet": r.get("body", ""),
                    })
    except Exception as e:
        log.warning(f"DDGS search failed: {e}")
        raise HTTPException(status_code=502, detail=f"DDGS error: {e}")

    return {
        "results": results,
        "count": len(results),
        "source": "ddgs",
        "query": req.query,
        "fetched_at": datetime.utcnow().isoformat(),
    }


# ---------------------------------------------------------------------------
# POST /api/search/news  — NewsData.io primary + DDGS fallback
# ---------------------------------------------------------------------------
NEWSDATA_API_KEY = os.getenv("NEWSDATA_API_KEY", "")

@app.post("/api/search/news")
def search_news(req: NewsRequest) -> dict[str, Any]:
    """
    News search pipeline:
      1. NewsData.io (primary — 2,000 articles/day free)
      2. DDGS news fallback if NewsData.io fails or key missing
    Returns list of {title, url, snippet, published, source}.
    """
    results: list[dict[str, Any]] = []
    source_used = "none"

    # --- Attempt 1: NewsData.io ---
    if NEWSDATA_API_KEY:
        try:
            url = "https://newsdata.io/api/1/latest"
            params = {
                "apikey": NEWSDATA_API_KEY,
                "q": req.query,
                "language": req.language,
                "size": min(req.max_results, 10),  # free tier max per call
            }
            resp = httpx.get(url, params=params, timeout=15)
            resp.raise_for_status()
            data = resp.json()
            articles = data.get("results", [])
            for a in articles:
                results.append({
                    "title": a.get("title", ""),
                    "url": a.get("link", ""),
                    "snippet": a.get("description", "") or a.get("content", "")[:300],
                    "published": a.get("pubDate", ""),
                    "source": a.get("source_id", ""),
                    "image_url": a.get("image_url"),
                })
            source_used = "newsdata.io"
            log.info(f"NewsData.io returned {len(results)} results for: {req.query}")
        except Exception as e:
            log.warning(f"NewsData.io failed: {e} — falling back to DDGS news")

    # --- Attempt 2: DDGS news fallback ---
    if not results:
        try:
            from duckduckgo_search import DDGS
            with DDGS() as ddgs:
                raw = ddgs.news(req.query, max_results=req.max_results)
                for r in raw:
                    results.append({
                        "title": r.get("title", ""),
                        "url": r.get("url", ""),
                        "snippet": r.get("body", ""),
                        "published": r.get("date", ""),
                        "source": r.get("source", ""),
                    })
            source_used = "ddgs_news"
            log.info(f"DDGS news returned {len(results)} results for: {req.query}")
        except Exception as e:
            log.warning(f"DDGS news also failed: {e}")

    return {
        "results": results,
        "count": len(results),
        "source": source_used,
        "query": req.query,
        "fetched_at": datetime.utcnow().isoformat(),
    }


# ---------------------------------------------------------------------------
# POST /api/search/extract-and-classify  — Full article text + risk classification data
# Fetches full text via Trafilatura, returns structured data for Claude to classify
# ---------------------------------------------------------------------------
class ArticleClassifyRequest(BaseModel):
    articles: list[dict[str, Any]]  # [{title, url, snippet, published, source}]
    supplier_name: str


@app.post("/api/search/extract-and-classify")
def extract_and_classify(req: ArticleClassifyRequest) -> dict[str, Any]:
    """
    For each article URL: extract full text via Trafilatura.
    Returns enriched articles with full_text and confidence.
    Claude does the actual classification on the Node side.
    """
    enriched: list[dict[str, Any]] = []

    for article in req.articles[:10]:  # cap at 10 per call
        url = article.get("url", "")
        enriched_article = dict(article)
        enriched_article["full_text"] = None
        enriched_article["extraction_confidence"] = "low"

        if url:
            try:
                downloaded = trafilatura.fetch_url(url)
                text = trafilatura.extract(downloaded) if downloaded else None
                if text and len(text) >= 200:
                    enriched_article["full_text"] = text[:8000]  # cap at 8k chars
                    enriched_article["extraction_confidence"] = "high"
                else:
                    # Try archive.ph
                    archive_url = f"https://archive.ph/{url}"
                    downloaded_arch = trafilatura.fetch_url(archive_url)
                    text_arch = trafilatura.extract(downloaded_arch) if downloaded_arch else None
                    if text_arch and len(text_arch) >= 200:
                        enriched_article["full_text"] = text_arch[:8000]
                        enriched_article["extraction_confidence"] = "medium"
            except Exception as e:
                log.warning(f"Article extraction failed for {url}: {e}")

        enriched.append(enriched_article)

    return {
        "articles": enriched,
        "supplier_name": req.supplier_name,
        "processed": len(enriched),
        "fetched_at": datetime.utcnow().isoformat(),
    }


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=5001, reload=True, log_level="info")
